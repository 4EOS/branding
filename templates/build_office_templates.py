#!/usr/bin/env python3
"""Generate 4EOS Office templates (Word + PowerPoint) from the brand spec.

Source of truth: ../4EOS-Brand-Styling-SKILL.md
Produces:
  4EOS-Document-Template.docx / .dotx   Word letterhead with brand styles
  4EOS-Presentation-Template.pptx / .potx  Title + content slide skeleton

Requires: python-docx, python-pptx, plus a PNG logo at the path in LOGO_PNG
(convert the official .webp once: python -c "from PIL import Image; Image.open('logo.webp').save('logo.png')").
Run from a venv: python3 -m venv env && env/bin/pip install python-docx python-pptx pillow
"""
import shutil
import zipfile
from pathlib import Path

from docx import Document
from docx.enum.text import WD_TAB_ALIGNMENT
from docx.oxml.ns import qn
from docx.oxml import OxmlElement
from docx.shared import Inches, Pt, RGBColor

from pptx import Presentation
from pptx.util import Inches as PInches, Pt as PPt, Emu
from pptx.dml.color import RGBColor as PRGBColor
from pptx.enum.text import PP_ALIGN

HERE = Path(__file__).resolve().parent
LOGO_PNG = Path("/tmp/4eos_logo.png")

# Brand tokens (4EOS-Brand-Styling-SKILL.md)
NAVY = "003D6C"
BLUE = "006699"
YELLOW = "DAAA00"
DARK = "1C2430"
MID = "4A6078"
TINT_H3 = "7EC8E3"
BRAND_LINE = "4EOS \u00b7 Computers / Networks / Security"


# ---------------------------------------------------------------- Word (docx)
def set_run_font(run, name, size_pt, color, bold=False, italic=False):
    run.font.name = name
    run.font.size = Pt(size_pt)
    run.font.color.rgb = RGBColor.from_string(color)
    run.font.bold = bold
    run.font.italic = italic
    rPr = run._element.get_or_add_rPr()
    rFonts = rPr.find(qn("w:rFonts"))
    if rFonts is None:
        rFonts = OxmlElement("w:rFonts")
        rPr.append(rFonts)
    for attr in ("w:ascii", "w:hAnsi", "w:cs"):
        rFonts.set(qn(attr), name)


def style_font(style, name, size_pt, color, bold=False):
    style.font.name = name
    style.font.size = Pt(size_pt)
    style.font.color.rgb = RGBColor.from_string(color)
    style.font.bold = bold
    rPr = style.element.get_or_add_rPr()
    rFonts = rPr.find(qn("w:rFonts"))
    if rFonts is None:
        rFonts = OxmlElement("w:rFonts")
        rPr.append(rFonts)
    for attr in ("w:ascii", "w:hAnsi", "w:cs"):
        rFonts.set(qn(attr), name)


def para_border(para, edge, color, size_eighths_pt, space="4"):
    """Add a single-edge border to a paragraph (size in eighths of a point)."""
    pPr = para._p.get_or_add_pPr()
    pBdr = OxmlElement("w:pBdr")
    el = OxmlElement(f"w:{edge}")
    el.set(qn("w:val"), "single")
    el.set(qn("w:sz"), str(size_eighths_pt))
    el.set(qn("w:space"), space)
    el.set(qn("w:color"), color)
    pBdr.append(el)
    pPr.append(pBdr)


def add_page_field(para):
    """Append a PAGE field to a paragraph."""
    fld = OxmlElement("w:fldSimple")
    fld.set(qn("w:instr"), "PAGE")
    run = OxmlElement("w:r")
    rPr = OxmlElement("w:rPr")
    rFonts = OxmlElement("w:rFonts")
    for attr in ("w:ascii", "w:hAnsi"):
        rFonts.set(qn(attr), "Calibri")
    rPr.append(rFonts)
    sz = OxmlElement("w:sz")
    sz.set(qn("w:val"), "20")  # 10pt
    rPr.append(sz)
    color = OxmlElement("w:color")
    color.set(qn("w:val"), MID)
    rPr.append(color)
    run.append(rPr)
    text = OxmlElement("w:t")
    text.text = "1"
    run.append(text)
    fld.append(run)
    para._p.append(fld)


def build_docx(path: Path):
    doc = Document()

    # Page: US Letter, 1in margins.
    section = doc.sections[0]
    section.page_width = Inches(8.5)
    section.page_height = Inches(11)
    section.top_margin = Inches(1.25)  # room for letterhead
    section.bottom_margin = Inches(1)
    section.left_margin = Inches(1)
    section.right_margin = Inches(1)

    # Core styles from the brand type scale (white-background column).
    styles = doc.styles
    normal = styles["Normal"]
    style_font(normal, "Calibri", 12, DARK)
    normal.paragraph_format.line_spacing = 1.15
    normal.paragraph_format.space_after = Pt(8)

    style_font(styles["Title"], "Bahnschrift", 28, NAVY, bold=True)
    style_font(styles["Heading 1"], "Bahnschrift", 26, NAVY, bold=True)
    style_font(styles["Heading 2"], "Bahnschrift", 20, BLUE, bold=True)
    style_font(styles["Heading 3"], "Bahnschrift", 15, NAVY, bold=True)
    style_font(styles["Caption"], "Calibri", 10, MID)
    styles["Caption"].font.italic = False

    # Header: logo left, navy rule under the letterhead row.
    header = section.header
    hp = header.paragraphs[0]
    run = hp.add_run()
    run.add_picture(str(LOGO_PNG), height=Inches(0.55))
    para_border(hp, "bottom", NAVY, 12, space="8")  # 1.5pt navy rule

    # Footer: brand line + page number, mid gray 10pt, thin yellow rule above.
    footer = section.footer
    fp = footer.paragraphs[0]
    fp.paragraph_format.tab_stops.add_tab_stop(Inches(6.5), WD_TAB_ALIGNMENT.RIGHT)
    r = fp.add_run(BRAND_LINE + "\t")
    set_run_font(r, "Calibri", 10, MID)
    r2 = fp.add_run("Page ")
    set_run_font(r2, "Calibri", 10, MID)
    add_page_field(fp)
    para_border(fp, "top", YELLOW, 8, space="4")  # 1pt yellow rule

    # Skeleton content: demonstrates every style; replace in real documents.
    doc.add_paragraph("Document title", style="Title")
    cap = doc.add_paragraph(
        "Subtitle or document summary. Calibri 10pt, mid gray.", style="Caption"
    )
    doc.add_paragraph("First-level heading", style="Heading 1")
    doc.add_paragraph(
        "Body text. Calibri 12pt in Dark Base on white, 1.15 line spacing. "
        "Replace this paragraph with real content."
    )
    doc.add_paragraph("Second-level heading", style="Heading 2")
    doc.add_paragraph("Body text under H2.")
    doc.add_paragraph("Third-level heading", style="Heading 3")
    doc.add_paragraph("Body text under H3.")
    doc.add_paragraph("First bullet point", style="List Bullet")
    doc.add_paragraph("Second bullet point", style="List Bullet")

    props = doc.core_properties
    props.title = "4EOS Document Template"
    props.author = "4EOS"
    props.comments = "Brand styles per 4EOS-Brand-Styling-SKILL.md"

    doc.save(path)

    # Theme fonts so the Office font picker defaults to Bahnschrift/Calibri.
    patch_theme_fonts(path, "word/theme/theme1.xml", "Bahnschrift", "Calibri")


# ----------------------------------------------------------- PowerPoint (pptx)
def ppt_textbox(slide, left, top, width, height):
    box = slide.shapes.add_textbox(left, top, width, height)
    box.text_frame.word_wrap = True
    return box


def set_para(para, text, font, size, color, bold=False, align=None):
    para.text = text
    if align is not None:
        para.alignment = align
    for run in para.runs:
        run.font.name = font
        run.font.size = PPt(size)
        run.font.bold = bold
        run.font.color.rgb = PRGBColor.from_string(color)


def build_pptx(path: Path):
    prs = Presentation()
    prs.slide_width = PInches(13.333)
    prs.slide_height = PInches(7.5)
    blank = prs.slide_layouts[6]

    # -- Title slide: navy field, yellow bottom bar, white Bahnschrift title.
    slide = prs.slides.add_slide(blank)
    bg = slide.shapes.add_shape(
        1, 0, 0, prs.slide_width, prs.slide_height  # 1 = rectangle
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = PRGBColor.from_string(NAVY)
    bg.line.fill.background()
    bg.shadow.inherit = False

    bar = slide.shapes.add_shape(1, 0, prs.slide_height - PInches(0.12), prs.slide_width, PInches(0.12))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRGBColor.from_string(YELLOW)
    bar.line.fill.background()
    bar.shadow.inherit = False

    title = ppt_textbox(slide, PInches(1), PInches(2.6), PInches(11.3), PInches(1.6))
    set_para(title.text_frame.paragraphs[0], "Presentation title", "Bahnschrift", 44, "FFFFFF", bold=True)
    sub = ppt_textbox(slide, PInches(1), PInches(4.1), PInches(11.3), PInches(0.9))
    set_para(sub.text_frame.paragraphs[0], "Subtitle, presenter, date", "Calibri", 20, TINT_H3)

    # -- Content slide: white, navy Bahnschrift title, Calibri body, logo top-right.
    slide = prs.slides.add_slide(blank)
    title = ppt_textbox(slide, PInches(0.8), PInches(0.5), PInches(9.6), PInches(1))
    set_para(title.text_frame.paragraphs[0], "Slide heading", "Bahnschrift", 28, NAVY, bold=True)
    slide.shapes.add_picture(
        str(LOGO_PNG), PInches(11.4), PInches(0.5), height=PInches(0.5)
    )
    # Yellow rule under the title row.
    rule = slide.shapes.add_shape(1, PInches(0.8), PInches(1.5), PInches(11.73), Emu(28575))  # ~2pt
    rule.fill.solid()
    rule.fill.fore_color.rgb = PRGBColor.from_string(YELLOW)
    rule.line.fill.background()
    rule.shadow.inherit = False

    body = ppt_textbox(slide, PInches(0.8), PInches(1.9), PInches(11.7), PInches(4.6))
    tf = body.text_frame
    set_para(tf.paragraphs[0], "First point. Calibri, Dark Base on white.", "Calibri", 18, DARK)
    for text in ("Second point.", "Third point."):
        p = tf.add_paragraph()
        set_para(p, text, "Calibri", 18, DARK)
    for p in tf.paragraphs:
        p.space_after = PPt(10)

    foot = ppt_textbox(slide, PInches(0.8), PInches(6.9), PInches(8), PInches(0.4))
    set_para(foot.text_frame.paragraphs[0], BRAND_LINE, "Calibri", 10, MID)

    props = prs.core_properties
    props.title = "4EOS Presentation Template"
    props.author = "4EOS"

    prs.save(path)
    patch_theme_fonts(path, "ppt/theme/theme1.xml", "Bahnschrift", "Calibri")


# ----------------------------------------------------------------- shared bits
def patch_theme_fonts(pkg: Path, theme_part: str, major: str, minor: str):
    """Rewrite theme1.xml major/minor latin typefaces inside an Office package."""
    tmp = pkg.with_suffix(".tmp")
    with zipfile.ZipFile(pkg, "r") as zin, zipfile.ZipFile(tmp, "w", zipfile.ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)
            if item.filename == theme_part:
                xml = data.decode("utf-8")
                xml = xml.replace('<a:latin typeface="Calibri Light"', f'<a:latin typeface="{major}"', 1)
                xml = xml.replace('<a:latin typeface="Calibri"', f'<a:latin typeface="{minor}"', 1)
                data = xml.encode("utf-8")
            zout.writestr(item, data)
    tmp.replace(pkg)


def derive_template(src: Path, dst: Path, doc_ct: str, tpl_ct: str):
    with zipfile.ZipFile(src, "r") as zin, zipfile.ZipFile(dst, "w", zipfile.ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)
            if item.filename == "[Content_Types].xml":
                data = data.decode("utf-8").replace(doc_ct, tpl_ct).encode("utf-8")
            zout.writestr(item, data)


def main():
    assert LOGO_PNG.exists(), f"logo PNG missing: {LOGO_PNG}"
    docx_path = HERE / "4EOS-Document-Template.docx"
    pptx_path = HERE / "4EOS-Presentation-Template.pptx"
    build_docx(docx_path)
    build_pptx(pptx_path)
    derive_template(
        docx_path,
        HERE / "4EOS-Document-Template.dotx",
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml",
        "application/vnd.openxmlformats-officedocument.wordprocessingml.template.main+xml",
    )
    derive_template(
        pptx_path,
        HERE / "4EOS-Presentation-Template.potx",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml",
        "application/vnd.openxmlformats-officedocument.presentationml.template.main+xml",
    )
    for f in sorted(HERE.glob("4EOS-*.*")):
        print(f.name, f.stat().st_size, "bytes")


if __name__ == "__main__":
    main()
